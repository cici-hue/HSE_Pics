import streamlit as st
import fitz  # PyMuPDF
# import PyMuPDF  # PyMuPDF
import re
import zipfile
import tempfile
import os
from pathlib import Path
from collections import defaultdict, OrderedDict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import hashlib
import io

# 设置页面配置
st.set_page_config(
    page_title="PDF缺陷图片提取与PPT生成系统",
    page_icon="📊",
    layout="wide"
)

class PDFDefectExtractor:
    """PDF缺陷提取器类"""
    def __init__(self):
        self.extracted_items = []
    
    def extract_defects_from_pdf(self, pdf_file, filename):
        """从单个PDF文件中提取缺陷信息"""
        extracted_items = []
        
        try:
            pdf_bytes = pdf_file.read()
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                blocks = page.get_text("dict")["blocks"]
                image_list = page.get_images(full=True)
                
                # 找出所有图片块
                image_blocks = []
                for i, block in enumerate(blocks):
                    if block["type"] == 1:  # 图片块
                        image_blocks.append({
                            "index": i,
                            "bbox": block["bbox"],
                            "y_position": block["bbox"][1]
                        })
                
                # 按y坐标排序（从上到下）
                image_blocks.sort(key=lambda x: x["y_position"])
                
                # 处理每个图片块（跳过第一个）
                for block_idx, block_info in enumerate(image_blocks):
                    if block_idx == 0:  # 跳过第一张图片
                        continue
                    
                    result = self._analyze_text_blocks(blocks, block_info["index"])
                    
                    if result and "defect_code" in result:
                        # 根据图片块的位置查找最接近的图片
                        bbox = block_info["bbox"]
                        matched_image_idx = self._find_matching_image(page, bbox, image_list)
                        
                        if matched_image_idx is not None:
                            try:
                                xref = image_list[matched_image_idx][0]
                                base_image = doc.extract_image(xref)
                                
                                # 清理缺陷原因作为文件名
                                reason = result.get("reason", f"defect_{result['defect_code']}")
                                clean_reason = self._sanitize_filename(reason)
                                
                                if not clean_reason or clean_reason == "_":
                                    clean_reason = f"defect_{result['defect_code']}"
                                
                                extracted_items.append({
                                    "pdf_name": filename,
                                    "page": page_num + 1,
                                    "defect_code": result.get("defect_code", ""),
                                    "reason": reason,
                                    "clean_reason": clean_reason,
                                    "image_data": base_image["image"],
                                    "image_ext": base_image["ext"]
                                })
                                
                            except Exception as e:
                                st.warning(f"提取图片失败: {e}")
                                continue
            
            doc.close()
        except Exception as e:
            st.error(f"处理PDF文件 {filename} 时出错: {str(e)}")
        
        return extracted_items
    
    def _analyze_text_blocks(self, blocks, start_index):
        """分析图片块后面的6个文本块"""
        result = {}
        text_blocks = []
        current_index = start_index + 1
        
        while len(text_blocks) < 6 and current_index < len(blocks):
            block = blocks[current_index]
            if block["type"] == 0:  # 文本块
                text = self._extract_text_from_block(block)
                if text.strip():
                    text_blocks.append(text)
            current_index += 1
        
        if len(text_blocks) < 6:
            return None
        
        # 检查第5个文本块
        if "defect code" not in text_blocks[4].lower():
            return None
        
        # 提取缺陷代码
        code_match = re.search(r'defect code\s+(\d+)', text_blocks[4], re.IGNORECASE)
        if not code_match:
            return None
        
        result["defect_code"] = code_match.group(1)
        
        # 提取原因
        if "defect" in text_blocks[5].lower():
            parts = re.split(r'\s+defect', text_blocks[5], flags=re.IGNORECASE)
            if parts and parts[0].strip():
                result["reason"] = parts[0].strip()
            else:
                return None
        else:
            return None
        
        return result
    
    def _find_matching_image(self, page, bbox, image_list):
        """查找匹配的图片"""
        block_center_x = (bbox[0] + bbox[2]) / 2
        block_center_y = (bbox[1] + bbox[3]) / 2
        
        best_match_idx = None
        min_distance = float('inf')
        
        for img_idx, img_info in enumerate(image_list):
            xref = img_info[0]
            img_rects = page.get_image_rects(xref)
            
            if img_rects:
                img_rect = img_rects[0]
                img_center_x = (img_rect.x0 + img_rect.x1) / 2
                img_center_y = (img_rect.y0 + img_rect.y1) / 2
                
                distance = ((img_center_x - block_center_x) ** 2 + 
                           (img_center_y - block_center_y) ** 2) ** 0.5
                
                if distance < min_distance:
                    min_distance = distance
                    best_match_idx = img_idx
        
        return best_match_idx
    
    def _extract_text_from_block(self, block):
        """从文本块中提取文本"""
        text = ""
        if "lines" in block:
            for line in block["lines"]:
                if "spans" in line:
                    for span in line["spans"]:
                        text += span.get("text", "") + " "
        return text.strip()
    
    def _sanitize_filename(self, filename):
        """清理文件名"""
        if not filename:
            return "unknown"
        
        # 移除特殊字符
        illegal_chars = ['<', '>', ':', '"', '/', '\\', '|', '?', '*', '\n', '\r', '\t']
        for char in illegal_chars:
            filename = filename.replace(char, '_')
        
        # 替换多个下划线为单个
        filename = re.sub(r'_{2,}', '_', filename)
        
        # 限制长度
        if len(filename) > 100:
            filename = filename[:100]
        
        return filename.strip()

class PPTCreator:
    """PPT生成器类"""
    def __init__(self):
        pass
    
    def create_ppt_from_images(self, all_defects, ppt_name="Defect_Report.pptx"):
        """从提取的图片创建PPT"""
        if not all_defects:
            return None
        
        # 按缺陷原因分类图片
        defects_by_reason = OrderedDict()
        file_counter = defaultdict(int)
        
        for defect in all_defects:
            reason = defect['reason']
            clean_reason = self._sanitize_filename(reason)
            
            if not clean_reason or clean_reason == "_":
                clean_reason = f"defect_{defect.get('defect_code', 'unknown')}"
            
            # 处理重复的文件名
            file_counter[clean_reason] += 1
            count = file_counter[clean_reason]
            
            if count > 1:
                clean_reason = f"{clean_reason}_{count}"
            
            if reason not in defects_by_reason:
                defects_by_reason[reason] = []
            
            defects_by_reason[reason].append({
                'order_number': defect.get('pdf_name', 'unknown').replace('.pdf', ''),
                'image_data': defect['image_data'],
                'image_ext': defect['image_ext'],
                'clean_name': clean_reason
            })
        
        # 按缺陷原因名称排序
        defects_by_reason = OrderedDict(sorted(defects_by_reason.items()))
        
        # 创建PPT
        return self._create_pptx_by_defect_reason(defects_by_reason, ppt_name)
    
    def _create_pptx_by_defect_reason(self, defects_by_reason, ppt_name):
        """创建基于缺陷原因分类的PPT"""
        try:
            # 创建PPT对象
            prs = Presentation()
            
            # 设置幻灯片尺寸（16:9）
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            
            # 添加标题页
            self._add_title_page(prs, len(defects_by_reason), 
                               sum(len(images) for images in defects_by_reason.values()))
            
            # 添加目录页
            self._add_table_of_contents(prs, defects_by_reason)
            
            # 为每种缺陷类型创建内容
            for defect_index, (defect_reason, images) in enumerate(defects_by_reason.items(), 1):
                # 添加缺陷类型标题页
                self._add_defect_title_page(prs, defect_reason, defect_index, len(defects_by_reason))
                
                # 将图片分组，每3张一组
                for i in range(0, len(images), 3):
                    img_group = images[i:i+3]
                    group_number = i // 3 + 1
                    total_groups = (len(images) - 1) // 3 + 1
                    
                    # 添加图片页
                    self._add_defect_images_page(prs, defect_reason, img_group, group_number, total_groups)
            
            # 添加结束页
            self._add_ending_page(prs)
            
            # 保存到内存
            ppt_buffer = io.BytesIO()
            prs.save(ppt_buffer)
            ppt_buffer.seek(0)
            
            return ppt_buffer
            
        except Exception as e:
            st.error(f"创建PPT失败: {str(e)}")
            return None
    
    def _add_title_page(self, prs, defect_types_count, total_images_count):
        """添加标题页"""
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Quality Defect Report"
        subtitle.text = f"By Defect Reason\n\n" \
                        f"Total Defect Types: {defect_types_count}\n" \
                        f"Total Images: {total_images_count}\n" \
                        f"Generated: {self._get_current_date()}"
        
        # 调整副标题字体大小
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.font.size = Pt(20)
    
    def _add_table_of_contents(self, prs, defects_by_reason):
        """添加目录页"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 添加标题
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(15)
        height = Inches(1)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = "Table of Contents"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        
        # 添加目录内容
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(14)
        height = Inches(6)
        
        content_box = slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_box.text_frame
        
        for i, (defect_reason, images) in enumerate(defects_by_reason.items(), 1):
            p = content_frame.add_paragraph()
            p.text = f"{i}. {defect_reason} ({len(images)} images)"
            p.font.size = Pt(20)
            p.level = 0
            p.space_after = Pt(5)
    
    def _add_defect_title_page(self, prs, defect_reason, defect_index, total_defects):
        """添加缺陷类型标题页"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 添加标题
        left = Inches(1)
        top = Inches(2)
        width = Inches(14)
        height = Inches(3)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        
        # 添加缺陷类型标题
        p = title_frame.paragraphs[0]
        p.text = defect_reason
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # 添加页码信息
        p = title_frame.add_paragraph()
        p.text = f"Defect Type {defect_index} of {total_defects}"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER
    
    def _add_defect_images_page(self, prs, defect_reason, img_group, group_number, total_groups):
        """添加缺陷图片页"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 添加缺陷原因标题（顶部）
        self._add_defect_header(slide, defect_reason, group_number, total_groups)
        
        # 添加图片
        self._add_images_with_order_numbers(slide, img_group)
        
        # 添加页码
        self._add_page_number(slide, group_number, total_groups)
    
    def _add_defect_header(self, slide, defect_reason, group_number, total_groups):
        """添加页眉：缺陷原因"""
        left = Inches(0.5)
        top = Inches(0.2)
        width = Inches(15)
        height = Inches(0.8)
        
        header_box = slide.shapes.add_textbox(left, top, width, height)
        header_frame = header_box.text_frame
        
        # 添加缺陷原因
        p = header_frame.paragraphs[0]
        p.text = f"Defect Reason: {defect_reason}"
        p.font.size = Pt(28)
        p.font.bold = True
        
        # 添加分组信息（如果需要）
        if total_groups > 1:
            p = header_frame.add_paragraph()
            p.text = f"Group {group_number} of {total_groups}"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(100, 100, 100)
    
    def _add_images_with_order_numbers(self, slide, img_group):
        """添加图片和订单号"""
        img_count = len(img_group)
        if img_count == 0:
            return
        
        # 根据图片数量设置不同的布局
        if img_count == 1:
            # 1张图片：居中显示
            width = Inches(8)
            height = Inches(5.38)
            left = (Inches(16) - width) / 2
            top = Inches(1.8)
            
            positions = [(left, top, width, height)]
            
        elif img_count == 2:
            # 2张图片：并排显示
            width = Inches(6)
            height = Inches(5.38)
            total_width = 2 * width + Inches(1)
            start_left = (Inches(16) - total_width) / 2
            top = Inches(1.8)
            
            positions = [
                (start_left, top, width, height),
                (start_left + width + Inches(1), top, width, height)
            ]
            
        else:  # img_count == 3
            # 3张图片：横向并排显示，使用新尺寸
            width = Inches(4.78)
            height = Inches(5.38)
            
            total_width = 3 * width + Inches(2 * 0.3)
            start_left = (Inches(16) - total_width) / 2
            top = Inches(1.8)
            
            positions = [
                (start_left, top, width, height),
                (start_left + width + Inches(0.3), top, width, height),
                (start_left + 2 * (width + Inches(0.3)), top, width, height)
            ]
        
        # 添加图片和订单号
        for i, (img_info, (left, top, width, height)) in enumerate(zip(img_group, positions)):
            try:
                # 保存图片到临时文件
                with tempfile.NamedTemporaryFile(delete=False, suffix=f".{img_info['image_ext']}") as tmp_file:
                    tmp_file.write(img_info['image_data'])
                    tmp_file_path = tmp_file.name
                
                # 添加订单号（在图片上方）
                self._add_order_number(slide, img_info['order_number'], left, top - Inches(0.4), width)
                
                # 添加图片
                slide.shapes.add_picture(tmp_file_path, left, top, width=width, height=height)
                
                # 删除临时文件
                os.unlink(tmp_file_path)
                
            except Exception as e:
                st.warning(f"添加图片失败: {e}")
    
    def _add_order_number(self, slide, order_number, left, top, width):
        """添加订单号标签"""
        height = Inches(0.3)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        text_frame.text = f"Order No: {order_number}"
        text_frame.paragraphs[0].font.size = Pt(20)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _add_page_number(self, slide, current_group, total_groups):
        """添加页码"""
        left = Inches(14.5)
        top = Inches(8.2)
        width = Inches(1)
        height = Inches(0.5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        text_frame.text = f"{current_group}/{total_groups}"
        text_frame.paragraphs[0].font.size = Pt(12)
        text_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
        text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    def _add_ending_page(self, prs):
        """添加结束页"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 添加结束语
        left = Inches(2)
        top = Inches(3)
        width = Inches(12)
        height = Inches(3)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        p = text_frame.paragraphs[0]
        p.text = "End of Report"
        p.font.size = Pt(36)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        p = text_frame.add_paragraph()
        p.text = "Quality Control Department"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER
    
    def _get_current_date(self):
        """获取当前日期"""
        from datetime import datetime
        return datetime.now().strftime("%Y-%m-%d")
    
    def _sanitize_filename(self, filename):
        """清理文件名"""
        if not filename:
            return "unknown"
        
        illegal_chars = ['<', '>', ':', '"', '/', '\\', '|', '?', '*', '\n', '\r', '\t']
        for char in illegal_chars:
            filename = filename.replace(char, '_')
        
        filename = re.sub(r'_{2,}', '_', filename)
        
        if len(filename) > 100:
            filename = filename[:100]
        
        return filename.strip()

def main():
    """主应用函数"""
    st.title("📊 PDF缺陷提取与PPT生成系统")
    st.markdown("""
    ### 功能说明：
    1. **上传PDF文件**：上传HSE包含缺陷图片的Claim report PDF格式文档
    2. **自动提取缺陷图片**：系统自动识别和提取缺陷图片
    3. **生成PPT报告**：自动生成按缺陷原因分类的PPT报告
    4. **下载结果**：可以下载提取的图片和生成的PPT
    """)
    
    # 创建两个主要功能选项卡
    tab1, tab2 = st.tabs(["📄 PDF缺陷提取", "📊 PPT生成"])
    
    with tab1:
        st.header("PDF缺陷图片提取")
        uploaded_files = st.file_uploader(
            "选择PDF文件（可多选）",
            type="pdf",
            accept_multiple_files=True,
            key="pdf_uploader"
        )
        
        if uploaded_files:
            extractor = PDFDefectExtractor()
            all_defects = []
            
            # 进度条
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            with st.spinner("正在处理PDF文件..."):
                for i, uploaded_file in enumerate(uploaded_files):
                    # 更新进度
                    progress = (i) / len(uploaded_files)
                    progress_bar.progress(progress)
                    status_text.text(f"正在处理: {uploaded_file.name} ({i+1}/{len(uploaded_files)})")
                    
                    # 提取缺陷
                    defects = extractor.extract_defects_from_pdf(uploaded_file, uploaded_file.name)
                    for defect in defects:
                        defect['pdf_file'] = uploaded_file.name
                        all_defects.append(defect)
                
                progress_bar.progress(1.0)
                status_text.text("处理完成!")
            
            if all_defects:
                st.success(f"✅ 提取完成! 共找到 {len(all_defects)} 个缺陷")
                
                # 显示统计信息
                col1, col2, col3 = st.columns(3)
                with col1:
                    st.metric("PDF文件数", len(uploaded_files))
                with col2:
                    st.metric("总缺陷数", len(all_defects))
                with col3:
                    # 统计缺陷类型
                    defect_types = len(set(d['reason'] for d in all_defects))
                    st.metric("缺陷类型数", defect_types)
                
                # 显示缺陷详情表格
                st.subheader("📋 缺陷详情")
                display_data = []
                for i, defect in enumerate(all_defects[:50], 1):  # 最多显示50条
                    display_data.append({
                        "序号": i,
                        "PDF文件": defect['pdf_name'],
                        "页码": defect['page'],
                        "缺陷代码": defect.get('defect_code', 'N/A'),
                        "缺陷原因": defect['reason']
                    })
                
                st.dataframe(display_data, use_container_width=True)
                
                # 创建ZIP文件供下载
                st.subheader("📥 下载提取的图片")
                
                with tempfile.TemporaryDirectory() as tmpdir:
                    # 创建ZIP文件
                    zip_buffer = io.BytesIO()
                    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                        # 按PDF文件创建文件夹
                        file_counter = defaultdict(int)
                        
                        for defect in all_defects:
                            pdf_name = Path(defect['pdf_name']).stem
                            reason = defect['clean_reason']
                            
                            # 处理重复的文件名
                            file_counter[(pdf_name, reason)] += 1
                            count = file_counter[(pdf_name, reason)]
                            
                            if count == 1:
                                filename = f"{reason}.{defect['image_ext']}"
                            else:
                                filename = f"{reason}_{count}.{defect['image_ext']}"
                            
                            # 完整的ZIP路径
                            zip_path = f"{pdf_name}/{filename}"
                            
                            # 添加到ZIP
                            zip_file.writestr(zip_path, defect['image_data'])
                    
                    # 创建下载按钮
                    zip_buffer.seek(0)
                    st.download_button(
                        label="📦 下载所有图片 (ZIP格式)",
                        data=zip_buffer,
                        file_name="extracted_defect_images.zip",
                        mime="application/zip",
                        help="点击下载包含所有提取图片的ZIP文件"
                    )
                
                # 预览部分图片
                st.subheader("🖼️ 图片预览")
                preview_cols = st.columns(4)
                
                for idx, defect in enumerate(all_defects[:8]):  # 最多预览8张
                    col_idx = idx % 4
                    with preview_cols[col_idx]:
                        # 显示图片
                        st.image(
                            defect['image_data'],
                            caption=f"{defect['reason']} (第{defect['page']}页)",
                            use_container_width=True
                        )
                
                # 保存提取结果到session state
                st.session_state.extracted_defects = all_defects
                st.success("✅ 提取结果已保存，可以切换到PPT生成标签页")
                
            else:
                st.warning("⚠️ 未找到任何缺陷信息")
    
    with tab2:
        st.header("PPT报告生成")
        
        if 'extracted_defects' not in st.session_state or not st.session_state.extracted_defects:
            st.info("👈 请先在左侧标签页上传并提取PDF文件")
        else:
            st.success(f"✅ 已加载 {len(st.session_state.extracted_defects)} 个缺陷")
            
            # PPT选项
            col1, col2 = st.columns(2)
            with col1:
                ppt_name = st.text_input("PPT文件名", "Defect_Report.pptx")
            with col2:
                ppt_layout = st.selectbox(
                    "PPT布局",
                    ["每页3张图片", "每页2张图片", "每页1张图片"],
                    index=0
                )
            
            # 生成PPT
            if st.button("🚀 生成PPT报告", type="primary"):
                with st.spinner("正在生成PPT..."):
                    ppt_creator = PPTCreator()
                    ppt_buffer = ppt_creator.create_ppt_from_images(
                        st.session_state.extracted_defects,
                        ppt_name
                    )
                
                if ppt_buffer:
                    st.success("✅ PPT生成成功!")
                    
                    # 下载按钮
                    st.download_button(
                        label="📥 下载PPT文件",
                        data=ppt_buffer,
                        file_name=ppt_name,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        help="点击下载生成的PPT报告"
                    )
                    
                    # 显示PPT统计信息
                    st.subheader("📊 PPT报告统计")
                    
                    # 统计缺陷类型
                    defects_by_reason = defaultdict(list)
                    for defect in st.session_state.extracted_defects:
                        defects_by_reason[defect['reason']].append(defect)
                    
                    stats_data = []
                    for reason, defects in sorted(defects_by_reason.items()):
                        stats_data.append({
                            "缺陷原因": reason,
                            "图片数量": len(defects),
                            "涉及的PDF文件": len(set(d['pdf_name'] for d in defects))
                        })
                    
                    st.dataframe(stats_data, use_container_width=True)
                else:
                    st.error("❌ PPT生成失败")

# 侧边栏信息
with st.sidebar:
    st.header("ℹ️ 使用说明")
    st.markdown("""
    ### 操作步骤：
    1. **上传HSE的claim report PDF文件**：
       - 点击"浏览文件"或拖放PDF文件
       - 支持多文件同时上传
    
    2. **提取缺陷图片**：
       - 系统自动识别PDF中的缺陷图片
       - 自动提取缺陷原因和代码
       - 生成图片预览和统计信息
    
    3. **生成PPT报告**：
       - 切换到PPT生成标签页
       - 设置PPT文件名和布局
       - 点击生成按钮创建PPT
    
    4. **下载结果**：
       - 下载提取的图片（ZIP格式）
       - 下载生成的PPT报告
    """)
    
    st.header("📈 系统信息")
    st.markdown("""
    - **版本**: 1.0.0
    - **更新日期**: 2024-01-20
    - **支持格式**: PDF文件
    - **输出格式**: JPEG图片 + PPT报告
    """)

if __name__ == "__main__":
    main()
